#!/usr/bin/env python3
"""Build the sequencing-benchmark defense deck from 立项答辩.pptx (v2, QA-fixed)."""
import copy
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image

TPL = "pptx_build/template.pptx"
OUT = "results/答辩_多代测序基准.pptx"
FIG = "results/figures"

# ---------- text helper: replace text, keep first run's formatting ----------
def set_text(tf, text):
    lines = text.split("\n")
    p0 = tf.paragraphs[0]
    if not p0.runs:
        p0.add_run()
    runs = p0.runs
    runs[0].text = lines[0]
    for r in runs[1:]:
        r._r.getparent().remove(r._r)
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    txbody = tf._txBody
    for ln in lines[1:]:
        newp = copy.deepcopy(p0._p)
        rs = newp.findall(qn('a:r'))
        for rr in rs[1:]:
            newp.remove(rr)
        rs = newp.findall(qn('a:r'))
        if rs:
            t = rs[0].find(qn('a:t'))
            if t is not None:
                t.text = ln
        txbody.append(newp)

def iter_tf(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_tf(sh.shapes)
        elif sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
        elif sh.has_text_frame:
            yield sh.text_frame

def apply_repl(slide, pairs):
    for tf in iter_tf(slide.shapes):
        cur = tf.text.strip()
        for old, new in pairs:
            if cur.startswith(old):
                set_text(tf, new)
                break

def remove_el(sh):
    sh._element.getparent().remove(sh._element)

def strip_pics(shapes):
    """Recursively remove all PICTURE and MEDIA shapes."""
    for sh in list(shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            strip_pics(sh.shapes)
        elif sh.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.MEDIA):
            remove_el(sh)

def remove_lines(slide):
    for sh in list(slide.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.LINE:
            remove_el(sh)

def set_title(slide, text, left=0.42, top=0.16, width=11.5):
    """Find the title text box (top-left) and set short text + widen box."""
    cand = None
    for sh in slide.shapes:
        if sh.has_text_frame and sh.top is not None and Emu(sh.top).inches < 1.0 \
                and Emu(sh.left).inches < 3.0 and sh.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            cand = sh
            break
    if cand is None:
        return
    set_text(cand.text_frame, text)
    cand.left = Inches(left)
    cand.width = Inches(8.6)          # ends ~9.0in, before the top-right logo (~9.3in)
    cand.text_frame.word_wrap = True
    for p in cand.text_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()
        pPr.set('algn', 'l')          # force left at XML level

def remove_empty_autoshapes(slide):
    for sh in list(slide.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                remove_el(sh)

def add_textbox(slide, lines, l, t, w, h, size=15, header=False, color="1E2761"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        r = p.add_run(); r.text = ln
        r.font.name = "Calibri"
        r.font.color.rgb = RGBColor.from_string(color)
        if header and i == 0:
            r.font.size = Pt(size + 5); r.font.bold = True
        else:
            r.font.size = Pt(size)
    return tb

def fit_pic(slide, path, l, t, w, h):
    """Place image centered within box (l,t,w,h) inches, preserving aspect."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = w / h
    if ar > box_ar:
        nw = w; nh = w / ar
    else:
        nh = h; nw = h * ar
    nl = l + (w - nw) / 2
    nt = t + (h - nh) / 2
    slide.shapes.add_picture(path, Inches(nl), Inches(nt), Inches(nw), Inches(nh))

prs = Presentation(TPL)
S = prs.slides

# ===================== SLIDE 1 — cover =====================
apply_repl(S[0], [
    ("基于视线追踪", "面向低算力场景的多代测序技术比较与小型基因组组装质量评估"),
    ("神经驱动的高拟真度", "基于模拟 Reads 的多代测序组装基准"),
    ("指导教师", "指导教师：______\n负责人：______\n组员：A、B、C、D"),
    ("时间：", "时间：2026.06"),
    ("H Y B R I D", "S E Q U E N C I N G   B E N C H M A R K   |   2 N D   v s   3 R D   G E N   |   L O W - R E S O U R C E"),
])

# ===================== SLIDE 2 — TOC =====================
apply_repl(S[1], [
    ("项目背景", "研究背景"),
    ("技术路线", "实验设计与结果"),
    ("计划预期", "结论与展望"),
])

# ===================== SLIDE 3 — PART 01 =====================
apply_repl(S[2], [("项目背景", "研究背景")])

# ===================== SLIDE 4 — 1.1 研究问题 =====================
strip_pics(S[3].shapes)
for sh in list(S[3].shapes):            # remove the left screenshot frame + red-circle annotation group
    if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
        remove_el(sh)
set_title(S[3], "1.1 研究问题")
apply_repl(S[3], [
    ("政策驱动", "核心问题：低算力 / 低成本下，如何选测序 + 组装？"),
    ("国家在", "DNA 测序像把一本书撕成小纸条再拼回。reads 拼回基因组的质量同时取决于「测序技术」（长度 / 错误率 / 深度）与「组装算法」。预算有限、只有普通电脑时，必须决定：选哪种技术、测多深、用哪个组装器——这往往凭经验而非受控比较。"),
])

# ===================== SLIDE 5 — 1.2 相关工作 =====================
strip_pics(S[4].shapes)
set_title(S[4], "1.2 相关工作与工具链")
apply_repl(S[4], [
    ("斯坦福大学", "统一开源工具链下的受控比较"),
    ("一种多功能的脑机接口", "本项目不自研算法，而在统一框架下组合成熟工具做受控比较。"),
    ("Kim, T.", "参考：ART(2012)·Badread(2019)·SPAdes(2012)·Flye(2019)·hifiasm(2021)·Unicycler(2017)·QUAST(2013)·BUSCO(2021)"),
    ("NOIR 的局限", "工具分类"),
    ("维度", "环节"), ("局限", "代表工具"),
    ("感知方式", "读长模拟"), ("屏幕闪烁刺激", "ART(Illumina) · Badread(长读长)"),
    ("控制范式", "短读组装"), ("离散指令触发", "SPAdes（de Bruijn 图）"),
    ("动作质量", "长读组装"), ("机械、分段式", "Flye · Raven · miniasm · hifiasm"),
    ("泛化能力", "评估"), ("依赖库中已有技能", "QUAST(有参)·BUSCO/Merqury(无参)·Prokka"),
])

# ===================== SLIDE 6 — 1.3 三个缺陷 =====================
strip_pics(S[5].shapes)
set_title(S[5], "1.3 已有比较的三个缺陷")
apply_repl(S[5], [
    ("控制维度低", "只报单一覆盖度：缺梯度，看不到饱和拐点与成本权衡。"),
    ("意图解码粗糙", "技术与工具绑定：每种技术只配一个组装器，混淆「数据」与「算法」效应。"),
    ("现有方法依赖", "只看连续性：忽略成本（$/Gb）与算力（内存 / 时间），无法支撑决策。"),
    ("动作缺乏泛化能力", "只看连续性：忽略成本（$/Gb）与算力（内存 / 时间），无法支撑决策。"),
])

# ===================== SLIDE 7 — PART 02 =====================
apply_repl(S[6], [("技术路线", "实验设计与结果")])

# ===================== SLIDE 8 — 2.1 实验矩阵 =====================
strip_pics(S[7].shapes)
set_title(S[7], "2.1 实验矩阵与方法")
apply_repl(S[7], [
    ("感知层", "数据"), ("眼动仪+EEG", "ART + Badread"), ("几何条件+语义条件", "4 技术 × 5 覆盖度"),
    ("生成层", "组装"), ("LatentHOI 扩散模型", "6 个组装器"), ("生成“想象中的”人手轨迹", "SPAdes / Flye / Raven\nminiasm / hifiasm / Unicycler"),
    ("执行层", "评估"), ("参数重定向", "QUAST · BUSCO"), ("驱动实体机械手", "成本 · 算力（内存 / 时间）"),
])
fit_pic(S[7], f"{FIG}/read_length_distribution.png", 1.14, 1.25, 11.05, 3.55)

# ===================== SLIDE 9 — 结果① 覆盖度饱和 =====================
strip_pics(S[8].shapes)
set_title(S[8], "2.2 结果①：覆盖度饱和 → 最小够用 20×")
apply_repl(S[8], [
    ("机械手在面对没见过的杯子", "短读长 50× 仍 ≥80 contigs（不闭合）；长读长 20× 即拼成单条 contig。"),
    ("核心模型", "关键现象"),
    ("LatentHOI（CVPR2025）", "Illumina contigs：1296→129→89→82→87（不闭合）\nONT R10 / HiFi / Hybrid：20× 起单 contig"),
    ("算法亮点", "结论"),
    ("生成式控制范式", "长读长在 20× 即闭合基因组，BUSCO 完整度随之达 100%。20× 是「最小够用覆盖度」。"),
])
fit_pic(S[8], f"{FIG}/coverage_genome_fraction_curve.png", 0.45, 1.5, 7.3, 4.6)

# ===================== SLIDE 10 — 结果② 技术 vs 工具 =====================
strip_pics(S[9].shapes)
set_title(S[9], "2.3 结果②：技术 vs 工具解耦")
apply_repl(S[9], [
    ("感知层", "测序技术"), ("眼动仪+EEG", "方差贡献 0.239"), ("几何条件+语义条件", "数据本身的影响"),
    ("生成层", "组装工具"), ("LatentHOI 扩散模型", "方差贡献 0.249"), ("生成“想象中的”人手轨迹", "算法选择的影响"),
    ("执行层", "结论"), ("参数重定向", "二者 ≈ 相等"), ("驱动实体机械手", "只说技术不说组装器不充分"),
])
fit_pic(S[9], f"{FIG}/n50_comparison.png", 1.14, 1.25, 11.05, 3.55)

# ===================== SLIDE 11 — 结果③ R9→R10 =====================
strip_pics(S[10].shapes)
set_title(S[10], "2.4 结果③：技术代次 R9 → R10")
apply_repl(S[10], [
    ("研究挑战", "做法"),
    ("解决思路", "发现"),
    ("我们将系统分为", "固定读长、只变准确率，隔离「代次」效应。"),
    ("空间定位瓶颈", "ONT R10（≈99%）在 20× 拼成单条、indel 近零；ONT R9（≈90%）在 30× 仍停 2 条且 indel 高得多（79.7 vs 0.0 /100kbp）。现代 ONT 已追平 HiFi 大半——「三代」内部并非铁板一块。"),
])

# ===================== SLIDE 12 — 结果④ 重复/质粒 + 组装图 =====================
strip_pics(S[11].shapes)
remove_lines(S[11])
for sh in list(S[11].shapes):
    if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
        remove_el(sh)
remove_empty_autoshapes(S[11])
set_title(S[11], "2.5 结果④：重复 / 质粒")
apply_repl(S[11], [
    ("视觉定位流", "短读长（SPAdes）：缠成一团"),
    ("神经意图", "长读长（ONT R10）：近线性"),
    ("目标：获取用户想操作对象", "rep_v1/v2/plasmid @30×：Illumina 80+ contigs；ONT R10 / HiFi 仅 1–3 contigs、近 100%。"),
    ("目标： 解析用户想要", "长读长能跨重复、回收小质粒——这是机制，不只是结果。"),
    ("解决思路", "组装图"),
    ("我们将系统分为", "同一基因组的组装图：左 = 短读长缠绕，右 = 长读长线性。"),
])
fit_pic(S[11], f"{FIG}/graph_illumina_30x_spades_orig_s13.png", 0.6, 3.3, 5.9, 3.8)
fit_pic(S[11], f"{FIG}/graph_ont_r10_30x_flye_orig_s13.png", 6.85, 4.7, 5.9, 1.5)
add_textbox(S[11], [
    "↑ 整条染色体拼成一条直线 = 1 条大 contig（+1 小片段）",
    "对比左侧短读长的「线团」，长读长几乎无分叉。",
], 6.9, 6.25, 5.9, 1.0, size=14, color="1E2761")

# ===================== SLIDE 13 — 结果⑤ 抛光 + 下游 =====================
strip_pics(S[12].shapes)
set_title(S[12], "2.6 结果⑤：GPU 抛光与下游")
for sh in list(S[12].shapes):  # drop the crowded 启示 block
    if sh.has_text_frame and sh.text_frame.text.strip().startswith(("技术核心", "语义约束")):
        remove_el(sh)
apply_repl(S[12], [
    ("研究挑战", "GPU 抛光（诚实结论）"),
    ("解决思路", "下游基因恢复"),
    ("我们如何将原本设想", "ONT 组装用 GPU（RTX 4090）跑 Medaka 抛光：模拟 reads 上几乎无改善、R9 反而略变差（mismatch 16.8→26.7）——因 Medaka 学的是真实 ONT 信号的错误分布；Medaka 与 Merqury 同因失效，揭示模拟基准的局限。"),
    ("借鉴 MotionCLIP", "Prokka：碎成 139 段的 Illumina 只恢复 14/22 个 rRNA（操纵子被坍缩）；长读长完整恢复 22/22。"),
    ("Tevet, G.", "评估：QUAST(有参) + BUSCO 55/55 + Merqury + Prokka(下游)"),
])
for sh in S[12].shapes:  # lift 下游 block up to fill the band
    if sh.has_text_frame:
        t = sh.text_frame.text.strip()
        if t.startswith("下游基因恢复") or t.startswith("Prokka"):
            sh.top = Inches(2.5)
fit_pic(S[12], f"{FIG}/polish_indel_comparison.png", 0.6, 3.55, 5.8, 3.3)
fit_pic(S[12], f"{FIG}/reference_free_audit.png", 6.9, 3.55, 5.8, 3.3)

# ===================== SLIDE 14 — 结果总结表 =====================
set_title(S[13], "2.7 结果总结")
apply_repl(S[13], [
    ("维度", "维度"), ("NOIR 2.0（Stanford）", "短读长（Illumina）"), ("本项目（Ours）", "长读长（ONT R10 / HiFi）"),
    ("感知方式", "连续性"), ("屏幕闪烁刺激", "50× 仍 ≥80 contigs，不闭合"), ("自然眼动", "20× 即单 contig"),
    ("控制范式", "准确性"), ("离散指令触发预设代码", "碱基准；难跨重复"), ("生成式模型实时规划", "HiFi 近完美，ONT 有 indel"),
    ("动作质量", "下游基因"), ("机械、分段式", "仅恢复 14/22 rRNA"), ("连续、拟人、平滑", "完整恢复 22/22"),
    ("泛化能力", "成本 / 算力"), ("依赖库中已有技能", "便宜但拼不完整"), ("可适应未见过的异形物体", "ONT R10@20× 仅 $0.92"),
])

# ===================== SLIDE 15 — PART 03 =====================
apply_repl(S[14], [("计划预期", "结论与展望")])

# ===================== SLIDE 16 — 3.1 结论与建议 =====================
strip_pics(S[15].shapes)
# remove advisor cards (nested groups w/ photo-fills + multi-frame text) BEFORE relabeling
for sh in list(S[15].shapes):
    if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
        txt = " ".join(tf.text for tf in iter_tf([sh]))
        if any(k in txt for k in ("龙霄潇", "房钰棋", "国自然", "医疗人工智能", "三维视觉")):
            remove_el(sh)
set_title(S[15], "3.1 结论与实用建议")
apply_repl(S[15], [
    ("指导教师", "核心结论"),
    ("硬件", "首选"), ("软件", "备选"), ("成员", "不推荐"),
    ("丰富算力资源", "推荐：ONT R10 @ 20× + Flye\n约 $0.92，全预算最优"),
    ("已完成LatentHOI", "预算更紧：ONT R10 可接受略碎；要最高准确率用 HiFi"),
    ("前期积累扎实", "纯短读长不适合需要完整基因（尤其 rRNA）的分析"),
])
# clean conclusion boxes (left), replacing the removed advisor cards
add_textbox(S[15], [
    "① 短读长拼不完整",
    "50× 仍 ≥80 contigs，丢失 rRNA 等重复区基因；不适合需要完整基因的下游分析。",
    "证据：Illumina 仅 14/22 rRNA、139 contigs",
], 0.55, 1.95, 5.5, 2.1, size=15, header=True)
add_textbox(S[15], [
    "② 长读长 20× 即闭合",
    "ONT R10 / HiFi / Hybrid 在 20× 拼成单条、近 100%；最小够用覆盖度 = 20×。",
    "证据：覆盖度饱和曲线 + Pareto 拐点",
], 0.55, 4.55, 5.5, 2.1, size=15, header=True)

# ===================== SLIDE 17 — 3.2 四人分工 =====================
strip_pics(S[16].shapes)
set_title(S[16], "3.2 四人分工与可复现性")
apply_repl(S[16], [("团队分工", "四人分工"), ("经费预算", "可复现性")])
# division cards: target each member's group so identical body text is not cross-matched
DIV = {
    "方丁龙": ("A · 数据与模拟矩阵", "参考变体、4 技术×5 覆盖度模拟、reads 统计"),
    "黄见源": ("B · 组装与工具解耦", "六组装器矩阵、资源记录、Bandage 组装图"),
    "陆姝一": ("C · 评估与分析", "QUAST + BUSCO、Pareto / 方差 / 重复 / 审计、全部图表"),
    "康智":  ("D · 工程化与 GPU", "环境、Medaka 抛光、Prokka、可复现 harness、统稿"),
}
for sh in S[16].shapes:
    if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
        tfs = [tf for tf in iter_tf([sh]) if tf.text.strip()]
        joined = " ".join(tf.text for tf in tfs)
        for name, (hdr, body) in DIV.items():
            if name in joined:
                for tf in tfs:
                    if name in tf.text:               # the name frame -> role header
                        set_text(tf, hdr)
                    elif tf.text.strip().startswith("负责"):  # the description frame
                        set_text(tf, body)
                break
# replace the misfit 12x5 budget table with a clean reproducibility panel
for sh in list(S[16].shapes):
    if sh.has_table:
        remove_el(sh)
add_textbox(S[16], [
    "配置驱动 —— config/matrix.yaml 定义全矩阵",
    "固定随机种子 —— seeds.yaml，结果可复现",
    "统一 run-id —— {tech}_{cov}x_{asm}_{ref}_s{seed}",
    "一键评估 —— QUAST + BUSCO 自动汇总主表",
    "分层设计 —— L1 主网格 + L3 重复/质粒难点",
    "算力开销 —— 41 个组装合计约 31 分钟（普通机器）",
], 7.25, 2.3, 5.5, 4.7, size=16)

# ===================== SLIDE 18 — 3.3 局限与展望 =====================
strip_pics(S[17].shapes)
set_title(S[17], "3.3 局限与未来工作")
apply_repl(S[17], [
    ("研究计划", "局限与方向"),
    ("前期", "局限①：仅单一小基因组\n→ 扩展更多高重复基因组"),
    ("中期", "局限②：仅模拟 reads\n→ 真实 ONT / HiFi 数据验证"),
    ("后期", "方向③：多种子置信区间\n+ 学习型策略推荐"),
    ("1月", "局限①"), ("2-7月", "局限②"), ("8-12月", "方向③"),
    ("预期成果", "交付物"),
    ("实物演示系统一套", "可复现基准（55 组装）· ECCV 格式论文 · 课程报告与答辩材料"),
])

# ===================== SLIDE 19 — quote =====================
apply_repl(S[18], [
    ("让机器不仅是工具", "不止于比较谁更好\n更要解释机制、量化权衡、给出可复现的决策"),
])

# ===================== SLIDE 20 — thanks (keep 恳请批评指导！) =====================

prs.save(OUT)
print("saved", OUT)
